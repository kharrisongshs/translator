from pptx import Presentation
from deep_translator import GoogleTranslator

# === SETTINGS ===
input_file = 'original.pptx'
output_file = 'translated.pptx'
source_lang = 'en'       # Change if source isn't English
target_lang = 'ar'       # Change to your desired language

# === TRANSLATOR SETUP ===
translator = GoogleTranslator(source=source_lang, target=target_lang)

# === LOAD PRESENTATION ===
prs = Presentation(input_file)

# === TRANSLATE SLIDE TEXT ===
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text.strip()
                    if original_text:
                        try:
                            translated = translator.translate(original_text)
                            run.text = translated
                        except Exception as e:
                            print(f"Translation failed for '{original_text}': {e}")

# === SAVE OUTPUT ===
prs.save(output_file)
print(f"✅ Translation complete! Saved as {output_file}")
