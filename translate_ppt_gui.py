import os
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from deep_translator import GoogleTranslator

# Supported languages
LANGUAGES = {
    "English": "en",
    "French": "fr",
    "Arabic": "ar",
    "Spanish": "es",
    "German": "de",
    "Chinese (Simplified)": "zh-CN",
    "Russian": "ru",
    "Japanese": "ja",
}

# Translation Function
def translate_pptx(file_path, source_lang, target_lang):
    try:
        prs = Presentation(file_path)
        translator = GoogleTranslator(source=source_lang, target=target_lang)

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text.strip()
                            if text:
                                try:
                                    run.text = translator.translate(text)
                                except Exception as e:
                                    print(f"Failed to translate '{text}': {e}")

        output_file = file_path.replace(".pptx", f"_{target_lang}.pptx")
        prs.save(output_file)
        return output_file

    except Exception as e:
        return None

# GUI Setup
def run_gui():
    def select_file():
        path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
        file_path.set(path)

    def run_translation():
        path = file_path.get()
        src = LANGUAGES[src_lang.get()]
        tgt = LANGUAGES[tgt_lang.get()]
        if not path or not os.path.isfile(path):
            messagebox.showerror("Error", "Please select a valid PowerPoint file.")
            return

        output = translate_pptx(path, src, tgt)
        if output:
            messagebox.showinfo("Success", f"Translation complete!\nSaved as:\n{output}")
        else:
            messagebox.showerror("Error", "Translation failed.")

    root = tk.Tk()
    root.title("PowerPoint Translator")
    root.geometry("400x320")
    root.resizable(False, False)

    file_path = tk.StringVar()
    src_lang = tk.StringVar(value="English")
    tgt_lang = tk.StringVar(value="Arabic")

    # UI Layout
    tk.Label(root, text="Select PowerPoint File:").pack(pady=5)
    tk.Entry(root, textvariable=file_path, width=40).pack(pady=5)
    tk.Button(root, text="Browse...", command=select_file).pack(pady=5)

    tk.Label(root, text="From Language:").pack(pady=5)
    tk.OptionMenu(root, src_lang, *LANGUAGES.keys()).pack()

    tk.Label(root, text="To Language:").pack(pady=5)
    tk.OptionMenu(root, tgt_lang, *LANGUAGES.keys()).pack()

    tk.Button(root, text="Translate", command=run_translation, bg="#4CAF50", fg="white", width=20).pack(pady=15)

    root.mainloop()

# Run it
if __name__ == "__main__":
    run_gui()
